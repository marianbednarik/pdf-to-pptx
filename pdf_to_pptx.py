import os
import sys
import shutil
from tqdm import tqdm, trange
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


class Converter(object):

    def __init__(self, pdf_file, dpi):
        self.pdf_file = pdf_file
        self.dpi = dpi
        self.pptx_file = pdf_file.replace('.pdf', '.pptx')
        self.total_pages = 1

    def split_pdf(self):

        try:
            print('Splitting pdf...')
            pages = convert_from_path(self.pdf_file, dpi=self.dpi)
            self.total_pages = len(pages)
        except Exception as ex:
            print('Failed to split pdf \n', ex)
            return

        try:
            print('Saving images...')
            os.mkdir('temp')
            pbar = trange(len(pages))
            for index, page in enumerate(pages):
                page.save(f'temp/page{index}.jpg', 'JPEG')
                pbar.update(1)
            pbar.close()
        except Exception as ex:
            print('Failed to save images\n', ex)
            if os.path.isdir('temp'):
                shutil.rmtree('temp')
            return

    def create_pptx(self):
        try:
            print('Creating pptx...')
            prs = Presentation('templates/16_9.pptx')
            print('Adding slides...')

            for slide_number in trange(self.total_pages):
                img_path = f'temp/page{slide_number}.jpg'
                new_slide = prs.slide_layouts[0]
                slide = prs.slides.add_slide(new_slide)
                slide.placeholders[1]
                title = slide.shapes.title
                height = Inches(7.5)
                title.text = f'Image {slide_number}'
                slide.shapes.add_picture(img_path, 0, 0, height=height)

            prs.save(self.pptx_file)
            if os.path.isdir('temp'):
                shutil.rmtree('temp')
            print(
                f'Successfully converted {self.pdf_file} to {self.pptx_file}')

        except Exception as ex:
            print('An error occured while creating presentation\n', ex)
            if os.path.isdir('temp'):
                shutil.rmtree('temp')
            return

    def convert(self):
        self.split_pdf()
        self.create_pptx()


if __name__ == '__main__':
    args = sys.argv[1:]
    Converter(pdf_file=args[0], dpi=args[1]
              if len(args) > 1 else 100).convert()
